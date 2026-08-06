import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation(output_path):
    prs = Presentation()
    # Set 16:9 Widescreen dimensions (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Cyberpunk Dark Theme
    COLOR_BG = RGBColor(15, 23, 42)         # #0f172a Deep Slate
    COLOR_CARD = RGBColor(30, 41, 59)       # #1e293b Dark Blue Slate
    COLOR_ACCENT = RGBColor(13, 245, 175)    # #0df5af Mint Green Accent
    COLOR_CYAN = RGBColor(0, 242, 254)      # #00f2fe Cyan Accent
    COLOR_TEXT_MAIN = RGBColor(248, 250, 252)# #f8fafc Bright White
    COLOR_TEXT_MUTED = RGBColor(148, 163, 184)# #94a3b8 Muted Grey
    COLOR_AMBER = RGBColor(252, 163, 17)    # #fca311 Amber Gold
    COLOR_CRIMSON = RGBColor(255, 42, 95)    # #ff2a5f Crimson Red

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="BRDS-PEC MAJOR PROJECT PRESENTATION"):
        # Dark Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()

        # Category Accent
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT

        # Title Text
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_MAIN

    # -------------------------------------------------------------------------
    # SLIDE 1: Title Slide
    # -------------------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG
    bg1.line.fill.background()

    # Title Card Accent Frame
    card1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.1))
    card1.fill.solid()
    card1.fill.fore_color.rgb = COLOR_CARD
    card1.line.color.rgb = COLOR_CYAN
    card1.line.width = Pt(1.5)

    title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.333), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "Behavioral Ransomware Detection System"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_MAIN

    p2 = tf.add_paragraph()
    p2.text = "Pre-Encryption Containment (BRDS-PEC) — Major Project Presentation"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_ACCENT
    p2.space_before = Pt(10)

    meta_box = slide1.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.333), Inches(1.5))
    tf_meta = meta_box.text_frame
    p_meta1 = tf_meta.paragraphs[0]
    p_meta1.text = "Focus: Deep LSTM Sequence Modeling • Real-Time Sysmon Telemetry • HMAC Arm Tokens • PyTorch Autograd XAI"
    p_meta1.font.size = Pt(12)
    p_meta1.font.color.rgb = COLOR_TEXT_MUTED

    p_meta2 = tf_meta.add_paragraph()
    p_meta2.text = "Evaluated Performance: 99.71% Validation Accuracy | 100% Recall | 25/25 Passing Unit Tests"
    p_meta2.font.size = Pt(12)
    p_meta2.font.bold = True
    p_meta2.font.color.rgb = COLOR_CYAN
    p_meta2.space_before = Pt(8)

    # -------------------------------------------------------------------------
    # SLIDE 2: Problem Statement & Vision
    # -------------------------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "1. Project Overview & Problem Statement")

    col1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    col1.fill.solid()
    col1.fill.fore_color.rgb = COLOR_CARD
    col1.line.color.rgb = COLOR_CRIMSON

    tf_c1 = col1.text_frame
    tf_c1.word_wrap = True
    p = tf_c1.paragraphs[0]
    p.text = "THE PROBLEM: Legacy Defenses Fail"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_CRIMSON

    bullets_c1 = [
        "Traditional Antivirus (AV) relies on static hash matching—bypassed by zero-day mutations.",
        "EDR solutions react AFTER files are encrypted or ransom notes appear.",
        "Ransomware moves rapidly: Shadow copy deletion occurs in under 10 seconds.",
        "Corporate network propagation happens via SMB/RDP within minutes."
    ]
    for b in bullets_c1:
        p_b = tf_c1.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(10.5)
        p_b.font.color.rgb = COLOR_TEXT_MAIN
        p_b.space_before = Pt(8)

    col2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    col2.fill.solid()
    col2.fill.fore_color.rgb = COLOR_CARD
    col2.line.color.rgb = COLOR_ACCENT

    tf_c2 = col2.text_frame
    tf_c2.word_wrap = True
    p = tf_c2.paragraphs[0]
    p.text = "OUR SOLUTION: Pre-Encryption Containment (PEC)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    bullets_c2 = [
        "Monitors early Sysmon behavioral sequences (5s sliding windows).",
        "Deep LSTM neural network predicts threat probability pre-encryption.",
        "Cryptographic HMAC-SHA256 alert verification and signed arm tokens.",
        "Instant host isolation (NIC disable, ARP/DNS flush, firewall block).",
        "Targeted process tree collapse with protected OS process denylist."
    ]
    for b in bullets_c2:
        p_b = tf_c2.add_paragraph()
        p_b.text = "• " + b
        p_b.font.size = Pt(10.5)
        p_b.font.color.rgb = COLOR_TEXT_MAIN
        p_b.space_before = Pt(8)

    # -------------------------------------------------------------------------
    # SLIDE 3: System Architecture & Data Flow
    # -------------------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "2. System Architecture & End-to-End Pipeline")

    arch_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    arch_box.fill.solid()
    arch_box.fill.fore_color.rgb = COLOR_CARD
    arch_box.line.color.rgb = COLOR_CYAN

    tf_arch = arch_box.text_frame
    tf_arch.word_wrap = True
    p = tf_arch.paragraphs[0]
    p.text = "End-to-End Execution Flow Architecture"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN

    flow_steps = [
        ("1. Endpoint Logging", "Microsoft Sysmon v15+ captures Event IDs 1 (Process), 3 (Network), 7 (Image), 11 (File Create), 12/13 (Registry), 23/26 (File Wipe)."),
        ("2. Temporal Aggregation", "pipeline/temporal_aggregator.py groups events into 5-second sliding UTC windows per host & process key."),
        ("3. Deep Neural Scoring", "ml_engine/lstm/infer.py evaluates 30-step sequences via Bidirectional LSTM with concatenated Mean + Max pooling."),
        ("4. Cryptographic Alert Signing", "containment/alert_integrity.py signs alerts with HMAC-SHA256 and issues single-use .arm_token files."),
        ("5. Host Isolation & Tree Collapse", "trigger_daemon.py validates token and invokes ContainHost.ps1 -Armed & kill_process_tree.ps1 -Armed."),
        ("6. SOC Dashboard & XAI", "Frontend UI displays rolling risk curve (Chart.js) and PyTorch autograd feature drivers (LSTMSHAPExplainer).")
    ]

    for title, desc in flow_steps:
        p_s = tf_arch.add_paragraph()
        p_s.text = f"{title}: {desc}"
        p_s.font.size = Pt(10)
        p_s.font.color.rgb = COLOR_TEXT_MAIN
        p_s.space_before = Pt(6)

    # -------------------------------------------------------------------------
    # SLIDE 4: Implementation Deep Dive
    # -------------------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "3. Implementation & System Engineering")

    # 3 Column Cards
    card_width = Inches(3.7)
    card_gap = Inches(0.3)

    # Card 1: ML Engine
    c1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), card_width, Inches(5.2))
    c1.fill.solid()
    c1.fill.fore_color.rgb = COLOR_CARD
    c1.line.color.rgb = COLOR_ACCENT
    tf_1 = c1.text_frame
    tf_1.word_wrap = True
    tf_1.paragraphs[0].text = "ML & Deep Learning Engine"
    tf_1.paragraphs[0].font.size = Pt(13)
    tf_1.paragraphs[0].font.bold = True
    tf_1.paragraphs[0].font.color.rgb = COLOR_ACCENT
    ml_points = [
        "PyTorch LSTM: 2 layers, hidden dim 64.",
        "Concatenated Mean + Max sequence pooling across 30 timesteps.",
        "RCE Protection: torch.load(weights_only=True) + SHA-256 manifest check.",
        "StandardScaler saved in .joblib sidecar.",
        "Feature Mean-Padding on short sequence inference."
    ]
    for pt in ml_points:
        p_pt = tf_1.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.size = Pt(9.5)
        p_pt.font.color.rgb = COLOR_TEXT_MAIN
        p_pt.space_before = Pt(6)

    # Card 2: Security & APIs
    c2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8) + card_width + card_gap, Inches(1.6), card_width, Inches(5.2))
    c2.fill.solid()
    c2.fill.fore_color.rgb = COLOR_CARD
    c2.line.color.rgb = COLOR_CYAN
    tf_2 = c2.text_frame
    tf_2.word_wrap = True
    tf_2.paragraphs[0].text = "Backend REST APIs & Security"
    tf_2.paragraphs[0].font.size = Pt(13)
    tf_2.paragraphs[0].font.bold = True
    tf_2.paragraphs[0].font.color.rgb = COLOR_CYAN
    api_points = [
        "Flask 2.3+ WSGI Blueprints.",
        "@require_api_key header auth via hmac.compare_digest constant-time check.",
        "SQL _safe_like() wildcard escaping (% and _).",
        "TelemetryWatchdog: raises SENSOR_SILENCED on >30s log gaps.",
        "CORS origin allowlist protection."
    ]
    for pt in api_points:
        p_pt = tf_2.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.size = Pt(9.5)
        p_pt.font.color.rgb = COLOR_TEXT_MAIN
        p_pt.space_before = Pt(6)

    # Card 3: Containment & UI
    c3 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8) + (card_width + card_gap)*2, Inches(1.6), card_width, Inches(5.2))
    c3.fill.solid()
    c3.fill.fore_color.rgb = COLOR_CARD
    c3.line.color.rgb = COLOR_AMBER
    tf_3 = c3.text_frame
    tf_3.word_wrap = True
    tf_3.paragraphs[0].text = "Host Isolation & SOC Dashboard"
    tf_3.paragraphs[0].font.size = Pt(13)
    tf_3.paragraphs[0].font.bold = True
    tf_3.paragraphs[0].font.color.rgb = COLOR_AMBER
    cont_points = [
        "HMAC-SHA256 alert digests & single-use .arm_token creation.",
        "ContainHost.ps1: NIC disable, ARP/DNS flush, Firewall block rules.",
        "kill_process_tree.ps1: Tree collapse with $PROTECTED_PROCESSES denylist.",
        "Dark-mode SOC Dashboard (Chart.js risk curve & live event feed).",
        "LSTMSHAPExplainer PyTorch autograd XAI modal."
    ]
    for pt in cont_points:
        p_pt = tf_3.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.size = Pt(9.5)
        p_pt.font.color.rgb = COLOR_TEXT_MAIN
        p_pt.space_before = Pt(6)

    # -------------------------------------------------------------------------
    # SLIDE 5: Testing Strategy & Verification
    # -------------------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "4. Software Testing & Quality Assurance")

    t_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    t_box.fill.solid()
    t_box.fill.fore_color.rgb = COLOR_CARD
    t_box.line.color.rgb = COLOR_ACCENT

    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    p = tf_t.paragraphs[0]
    p.text = "PyTest Automated Test Suite (25 / 25 Passing Unit Tests)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    test_modules = [
        ("tests/test_backend.py", "Verifies API key authentication, live telemetry scoring, incident creation, and CORS allowlists."),
        ("tests/test_containment.py", "Validates HMAC-SHA256 signature checks, arm token creation, and dry-run PowerShell execution."),
        ("tests/test_database.py", "Tests SQLite persistence, SQLAlchemy ORM queries, and _safe_like() SQL wildcard escaping."),
        ("tests/test_lstm.py & test_lstm_integration.py", "Verifies PyTorch LSTM sequence dimensions, Mean+Max pooling, weights_only loading, and SHA-256 hash checks."),
        ("tests/test_ml_engine.py", "Tests source-level train/validation dataset splits, Isolation Forest scoring, and lead-time calculation helpers."),
        ("tests/test_pipeline.py", "Validates Sysmon XML parsing, 5s temporal windowing, vectorization, and TelemetryWatchdog gap detection."),
        ("tests/test_xai.py", "Tests LSTMSHAPExplainer PyTorch autograd feature attributions and sanitized API error response JSON.")
    ]

    for file_name, desc in test_modules:
        p_tm = tf_t.add_paragraph()
        p_tm.text = f"✔ {file_name}: {desc}"
        p_tm.font.size = Pt(9.5)
        p_tm.font.color.rgb = COLOR_TEXT_MAIN
        p_tm.space_before = Pt(5)

    # -------------------------------------------------------------------------
    # SLIDE 6: Results & Model Performance Metrics
    # -------------------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "5. Evaluation Results & Performance Metrics")

    # 4 Metric Callout Boxes
    box_w = Inches(2.7)
    box_h = Inches(2.0)

    metrics = [
        ("99.71%", "Validation Accuracy", "PyTorch LSTM Sequence Model", COLOR_ACCENT),
        ("100.0%", "Detection Recall", "Zero False Negatives on Attack Windows", COLOR_CYAN),
        ("98.69%", "Precision Score", "High Quality Alert Generation", COLOR_AMBER),
        ("0.988", "ROC-AUC Score", "Excellent Binary Classification", COLOR_ACCENT)
    ]

    for idx, (val, label, sub, col) in enumerate(metrics):
        x_pos = Inches(0.8) + idx * Inches(2.95)
        m_card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(1.6), box_w, box_h)
        m_card.fill.solid()
        m_card.fill.fore_color.rgb = COLOR_CARD
        m_card.line.color.rgb = col

        tf_m = m_card.text_frame
        tf_m.word_wrap = True
        p_v = tf_m.paragraphs[0]
        p_v.text = val
        p_v.font.size = Pt(28)
        p_v.font.bold = True
        p_v.font.color.rgb = col
        p_v.alignment = PP_ALIGN.CENTER

        p_l = tf_m.add_paragraph()
        p_l.text = label
        p_l.font.size = Pt(11)
        p_l.font.bold = True
        p_l.font.color.rgb = COLOR_TEXT_MAIN
        p_l.alignment = PP_ALIGN.CENTER

        p_s = tf_m.add_paragraph()
        p_s.text = sub
        p_s.font.size = Pt(8.5)
        p_s.font.color.rgb = COLOR_TEXT_MUTED
        p_s.alignment = PP_ALIGN.CENTER

    # Summary Table Below
    res_summary = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.9), Inches(11.7), Inches(2.9))
    res_summary.fill.solid()
    res_summary.fill.fore_color.rgb = COLOR_CARD
    res_summary.line.color.rgb = COLOR_CYAN

    tf_rs = res_summary.text_frame
    tf_rs.word_wrap = True
    tf_rs.paragraphs[0].text = "Key Experimental Takeaways & Reaction Speed"
    tf_rs.paragraphs[0].font.size = Pt(13)
    tf_rs.paragraphs[0].font.bold = True
    tf_rs.paragraphs[0].font.color.rgb = COLOR_CYAN

    takeaways = [
        "Pre-Encryption Reaction Time: 5 to 15 seconds after initial ransomware process execution.",
        "False Positive Rate (FPR): 5.75% on initial benchmark split (further reduced via real endpoint baselining).",
        "Source-Level Dataset Split: Evaluated using strict source-level splits to prevent scenario leakage between train and test sets.",
        "Attack Coverage: Evaluated against WannaCry, LockBit, Ryuk, and Sodinokibi execution telemetry."
    ]
    for t_kw in takeaways:
        p_tk = tf_rs.add_paragraph()
        p_tk.text = "• " + t_kw
        p_tk.font.size = Pt(10)
        p_tk.font.color.rgb = COLOR_TEXT_MAIN
        p_tk.space_before = Pt(6)

    # -------------------------------------------------------------------------
    # SLIDE 7: Cost Estimation & Resource Analysis
    # -------------------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "6. Cost Estimation & Resource Requirements")

    c_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    c_box.fill.solid()
    c_box.fill.fore_color.rgb = COLOR_CARD
    c_box.line.color.rgb = COLOR_AMBER

    tf_c = c_box.text_frame
    tf_c.word_wrap = True
    tf_c.paragraphs[0].text = "Resource & Cost Breakdown for Academic & Prototype Deployment"
    tf_c.paragraphs[0].font.size = Pt(14)
    tf_c.paragraphs[0].font.bold = True
    tf_c.paragraphs[0].font.color.rgb = COLOR_AMBER

    cost_items = [
        ("Software & Licensing Costs", "$0 (Zero Dollars)", "Built entirely using open-source tools: Python 3.14, PyTorch, Flask, Scikit-Learn, Microsoft Sysmon, PowerShell, and Chart.js."),
        ("Endpoint CPU / RAM Overhead", "< 1.5% CPU | ~45 MB RAM", "Lightweight Sysmon XML parsing and 5s temporal windowing impose minimal overhead on local endpoint hardware."),
        ("Model Training Infrastructure", "$0 (Standard Local PC)", "Trained locally in under 3 minutes on CPU/GPU without requiring expensive cloud AI server clusters."),
        ("Storage Footprint", "~15 MB per Endpoint / Day", "Compressed Sysmon event telemetry logs require minimal storage under standard rotation policies."),
        ("Enterprise ROI Comparison", "Savings of $1.5M+ per Breach", "The average enterprise ransomware breach costs $4.5M in downtime and ransom. BRDS prevents mass encryption at zero software cost.")
    ]

    for item_title, cost_val, detail in cost_items:
        p_ci = tf_c.add_paragraph()
        p_ci.text = f"• {item_title} [{cost_val}]: {detail}"
        p_ci.font.size = Pt(10)
        p_ci.font.color.rgb = COLOR_TEXT_MAIN
        p_ci.space_before = Pt(8)

    # -------------------------------------------------------------------------
    # SLIDE 8: Conclusion & Future Enhancements
    # -------------------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "7. Conclusion & Future Enhancements")

    col_f1 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    col_f1.fill.solid()
    col_f1.fill.fore_color.rgb = COLOR_CARD
    col_f1.line.color.rgb = COLOR_ACCENT
    tf_f1 = col_f1.text_frame
    tf_f1.word_wrap = True
    tf_f1.paragraphs[0].text = "Project Conclusion"
    tf_f1.paragraphs[0].font.size = Pt(14)
    tf_f1.paragraphs[0].font.bold = True
    tf_f1.paragraphs[0].font.color.rgb = COLOR_ACCENT

    conc_points = [
        "Successfully developed and validated a Pre-Encryption Containment (PEC) prototype.",
        "Deep LSTM model achieves 99.71% validation accuracy and 100% recall on ransomware attack windows.",
        "Cryptographic HMAC signing and protected process denylists ensure safe automated containment.",
        "All 25 automated unit tests pass cleanly, verifying backend, ML engine, and containment security."
    ]
    for pt in conc_points:
        p_pt = tf_f1.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.size = Pt(10.5)
        p_pt.font.color.rgb = COLOR_TEXT_MAIN
        p_pt.space_before = Pt(8)

    col_f2 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.2))
    col_f2.fill.solid()
    col_f2.fill.fore_color.rgb = COLOR_CARD
    col_f2.line.color.rgb = COLOR_CYAN
    tf_f2 = col_f2.text_frame
    tf_f2.word_wrap = True
    tf_f2.paragraphs[0].text = "Future Enhancements & Roadmap"
    tf_f2.paragraphs[0].font.size = Pt(14)
    tf_f2.paragraphs[0].font.bold = True
    tf_f2.paragraphs[0].font.color.rgb = COLOR_CYAN

    fut_points = [
        "Windows Kernel Driver Packaging: Move containment to kernel space for sub-millisecond execution.",
        "Central SIEM Connectors: Integrate REST APIs directly into Splunk, Microsoft Sentinel, and Elastic Security.",
        "Asynchronous Task Queue: Implement Celery + Redis for high-throughput enterprise event streams (>10k ev/s).",
        "Adaptive Online Retraining: Continuous model learning from enterprise background baseline drift."
    ]
    for pt in fut_points:
        p_pt = tf_f2.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.size = Pt(10.5)
        p_pt.font.color.rgb = COLOR_TEXT_MAIN
        p_pt.space_before = Pt(8)

    prs.save(output_path)
    print(f"Successfully generated PowerPoint presentation at: {output_path}")

if __name__ == '__main__':
    target = sys.argv[1] if len(sys.argv) > 1 else "docs/BRDS_Project_Presentation.pptx"
    create_presentation(target)
